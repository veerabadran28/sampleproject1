import streamlit as st
import boto3
from anthropic import Anthropic
from botocore.config import Config
import shutil
import os
import pandas as pd
import time
import json
import base64
import io
import re
import numpy as np
import openpyxl
from python_calamine import CalamineWorkbook
from openpyxl.cell import Cell
from openpyxl.worksheet.cell_range import CellRange
from boto3.dynamodb.conditions import Key 
import uuid
from pptx import Presentation
from botocore.exceptions import ClientError
from textractor import Textractor
from textractor.visualizers.entitylist import EntityList
from textractor.data.constants import TextractFeatures
from textractor.data.text_linearization_config import TextLinearizationConfig
import pytesseract
from PIL import Image
import PyPDF2
import chardet
from datetime import datetime    
from docx import Document as DocxDocument
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.document import Document
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from docx.table import Table as DocxTable
import concurrent.futures
from functools import partial
import csv
import textract
import requests
from streamlit_javascript import st_javascript
from urllib.parse import urljoin


#for document summarization
import tiktoken
from langchain.text_splitter import CharacterTextSplitter, RecursiveCharacterTextSplitter
from PyPDF2 import PdfReader


#for FAISS in-memory vectorembeddings
from langchain.embeddings import BedrockEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain_community.chat_models import BedrockChat
from langchain_core.prompts import ChatPromptTemplate
from langchain.embeddings import BedrockEmbeddings
from langchain.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import AmazonTextractPDFLoader
from langchain.chains import ConversationalRetrievalChain

from logging import getLogger
log = getLogger("test")
log.error("this is a test")  

config = Config(
    read_timeout=600, # Read timeout parameter
    retries = dict(
        max_attempts = 10 ## Handle retries
    )
)

# st.set_page_config(initial_sidebar_state="auto")
st.set_page_config(page_title="ESG Data Assistant 2", layout="wide")
# Read credentials
with open('config.json','r',encoding='utf-8') as f:
    config_file = json.load(f)
# pricing info
with open('pricing.json','r',encoding='utf-8') as f:
    pricing_file = json.load(f)

S3 = boto3.client('s3')
DYNAMODB  = boto3.resource('dynamodb')

LOCAL_CHAT_FILE_NAME = "chatbot-history.json"
DYNAMODB_TABLE=config_file["DynamodbTable"]
BUCKET=config_file["Bucket_Name"]
OUTPUT_TOKEN=config_file["max-output-token"]
S3_DOC_CACHE_PATH=config_file["document-upload-cache-s3-path"]
TEXTRACT_RESULT_CACHE_PATH=config_file["AmazonTextract-result-cache"]
LOAD_DOC_IN_ALL_CHAT_CONVO=config_file["load-doc-in-chat-history"]
CHAT_HISTORY_LENGTH=config_file["chat-history-loaded-length"]
DYNAMODB_USER=config_file["UserId"]
REGION=config_file["bedrock-region"]
USE_TEXTRACT=config_file["AmazonTextract"]
CSV_SEPERATOR=config_file["csv-delimiter"]

#Switch Buttons
DOCUMENT_SUMMARIZATION=True
QNA=False




bedrock_runtime = boto3.client(service_name='bedrock-runtime',region_name=REGION,config=config)

if 'messages' not in st.session_state:
    st.session_state['messages'] = []
if 'input_token' not in st.session_state:
    st.session_state['input_token'] = 0
if 'output_token' not in st.session_state:
    st.session_state['output_token'] = 0
if 'chat_hist' not in st.session_state:
    st.session_state['chat_hist'] = []
if 'user_sess' not in st.session_state:
    st.session_state['user_sess'] =str(time.time())
if 'chat_session_list' not in st.session_state:
    st.session_state['chat_session_list'] = []
if 'count' not in st.session_state:
    st.session_state['count'] = 0
if 'userid' not in st.session_state:
    st.session_state['userid']= config_file["UserId"]
if 'cost' not in st.session_state:
    st.session_state['cost'] = 0

#Get the user details from authentication
# config_file["UserId"] = "sam@sam.com"

def create_prompt_template():
    _template = """
        Given the following chat history and a follow up question, rephrase the follow up question to be standalone questions, in english. Skip the preamble and just get to the question.
        
        <chat_history>{chat_history}</chat_history>
        <follow_up_question>{question}</follow_up_question>
        
    """
    conversation_prompt = ChatPromptTemplate.from_template(_template)
    return conversation_prompt



def InMemoryFAISS(result_string, chat_history, question) -> str:

    bedrock = boto3.client('bedrock-runtime',region_name='us-east-1')
    bedrock_llm = BedrockChat(model_id="anthropic.claude-3-sonnet-20240229-v1:0", model_kwargs={"max_tokens": 3000,
                                                                                            "temperature": 0.0
                                                                                            })

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000,
                                                separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""],
                                                chunk_overlap=100)
    texts = text_splitter.split_text(result_string)
    embeddings = BedrockEmbeddings(client=bedrock,model_id="amazon.titan-embed-text-v2:0")
    vector_db = FAISS.from_texts(texts, embeddings)

    template = """
        Answer the question as truthfully as possible strictly using only the provided text, and if the answer is not contained within the text, say "I don't know". Skip any preamble text and reasoning and give just the answer. If user greets you, just greet back.

            <text>
            {context}
            {chat_history}
            </text>

            <question>
            {question}
            </question>
            <answer>

            """
    qa_prompt = ChatPromptTemplate.from_template(template)
    retreiver = vector_db.as_retriever(search_type='mmr', search_kwargs={'k': 5})
    qa = ConversationalRetrievalChain.from_llm(llm=bedrock_llm,
                                          retriever=retreiver,
                                           condense_question_prompt=create_prompt_template(),
                                           condense_question_llm=BedrockChat(model_id="anthropic.claude-3-sonnet-20240229-v1:0"),
                                           combine_docs_chain_kwargs={"prompt": qa_prompt},
                                           # verbose=True
                                          )
    result = qa({"question": question,
                "chat_history": chat_history,
                })
    return result["answer"]

    


def summarizer(prompt_data) -> str:
    """
    This function creates the summary of each individual chunk as well as the final summary.
    :param prompt_data: This is the prompt along with the respective chunk of text, at the end it contains all summary chunks combined.
    :return: A summary of the respective chunk of data passed in or the final summary that is a summary of all summary chunks.
    """
    # setting the key parameters to invoke Amazon Bedrock
    print(f"messages: {prompt_data}")
    body = json.dumps({"messages": prompt_data,
                       "max_tokens": 8191,
                       "temperature": 0,
                       "top_k": 250,
                       "top_p": 0.5,
                       "stop_sequences": [],
                        "anthropic_version": "bedrock-2023-05-31",
                       })
    # the specific Amazon Bedrock model you are using
    modelId = 'anthropic.claude-3-sonnet-20240229-v1:0'
    # type of data that should be expected upon invocation
    accept = 'application/json'
    contentType = 'application/json'
    # the invocation of bedrock, with all of the parameters you have configured
    # response = bedrock_runtime.invoke_model(body=body,
    #                                 modelId=modelId,
    #                                 accept=accept,
    #                                 contentType=contentType)
    
    response = bedrock_runtime.invoke_model(body=body,
                                                                 modelId=modelId,
                                                                 accept=accept,
                                                                 contentType=contentType
                                                                 )
    

    # gathering the response from bedrock, and parsing to get specifically the answer
    response_body = json.loads(response.get('body').read())
    answer = response_body["content"][0]["text"]
    # answer = response_body.get('completion')
    print(f"answer: {answer}")
    # returning the final summary for that chunk of text
    return answer


def num_tokens_from_string(string) -> int:
    """Returns the number of tokens in a text string."""
    encoding = tiktoken.get_encoding("cl100k_base")
    num_tokens = len(encoding.encode(string))
    return num_tokens


def Chunk_and_Summarize(text_from_textract) -> str:

    text = text_from_textract
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=100000,
        chunk_overlap=10000,
        length_function=len,
        add_start_index=True
    )
    # using the text splitter to split the entire string of text that contains all the text content of our PDF
    texts = text_splitter.create_documents([text])


    # Creating an empty summary string, as this is where we will append the summary of each chunk
    summary = ""
    # looping through each chunk of text we created, passing that into our prompt and generating a summary of that chunk
    for index, chunk in enumerate(texts):
        # gathering the text content of that specific chunk
        chunk_content = chunk.page_content
        # creating the prompt that will be passed into Bedrock with the text content of the chunk
        prompt = f"""\n\nHuman: Provide a detailed summary for the chunk of text provided to you:
        Text: {chunk_content}
        \n\nAssistant:"""

        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": f"Provide a detailed summary for the chunk of text provided to you: {chunk_content}"
                    }
                ]
            }
        ]

        # passing the prompt into the summarizer function to generate the summary of that chunk, and appending it to
        # the summary string
        summary+= summarizer(messages)
        # printing out the number of tokens contained in each chunk to provide a status update
        print(f"\n\nNumber of tokens for Chunk {index + 1} with the prompt: {num_tokens_from_string(prompt)} tokens")
        print("-------------------------------------------------------------------------------------------------------")
    # after we have generated the summaries of each chunk of text, and appended them to the single summary string,
    # we pass it into the final summary prompt
    final_summary_prompt = f"""\n\nHuman: You will be given a set of summaries from a document. Create a cohesive 
    summary from the provided individual summaries. The summary should very detailed and at least 10 pages. 
    Summaries: {summary}
            \n\nAssistant:"""
    
    print(summary)
    final_messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": f"You will be given a set of summaries from a document. Create a cohesive summary from the provided individual summaries. The summary should very detailed and at least 10 pages. Summaries: {summary}"
                }
            ]
        },
        {
            "role": "assistant",
            "content": ""
        }
        ]

    # print the total number of tokens being passed into the final summarization prompt
    print(f"Number of tokens for this Chunk with the final prompt: {num_tokens_from_string(final_summary_prompt)}")
    # generating the final summary of all the summaries we have previously generated.
    return summarizer(final_messages)



# def Chunk_and_Summarize(uploaded_file) -> str:
#     """
#     This function takes in the path to the file that was just uploaded through the streamlit app.
#     :param uploaded_file: This is a file path, that should point to the newly uploaded file that is temporarily stored
#     within the directory of this application.
#     :return: This returns the final summary of the PDF document that was initially passed in by the user through the
#     streamlit app.
#     """
#     # using PyPDF PdfReader to read in the PDF file as text
#     reader = PdfReader(uploaded_file)
#     # creating an empty string for us to append all the text extracted from the PDF
#     text = ""
#     # a simple for loop to iterate through all pages of the PDF we uploaded
#     for page in reader.pages:
#         # as we loop through each page, we extract the text from the page and append it to the "text" string
#         text += page.extract_text() + "\n"
#     # creating the text splitter, we are specifically using the the recursive text splitter from langchain:
#     # https://python.langchain.com/docs/modules/data_connection/document_transformers/text_splitters/recursive_text_splitter
    
    
    
#     text_splitter = RecursiveCharacterTextSplitter(
#         chunk_size=1000,
#         chunk_overlap=100,
#         length_function=len,
#         add_start_index=True
#     )
#     # using the text splitter to split the entire string of text that contains all the text content of our PDF
#     texts = text_splitter.create_documents([text])


#     # Creating an empty summary string, as this is where we will append the summary of each chunk
#     summary = ""
#     # looping through each chunk of text we created, passing that into our prompt and generating a summary of that chunk
#     for index, chunk in enumerate(texts):
#         # gathering the text content of that specific chunk
#         chunk_content = chunk.page_content
#         # creating the prompt that will be passed into Bedrock with the text content of the chunk
#         prompt = f"""\n\nHuman: Provide a detailed summary for the chunk of text provided to you:
#         Text: {chunk_content}
#         \n\nAssistant:"""
#         # passing the prompt into the summarizer function to generate the summary of that chunk, and appending it to
#         # the summary string
#         summary += summarizer(prompt)
#         # printing out the number of tokens contained in each chunk to provide a status update
#         print(f"\n\nNumber of tokens for Chunk {index + 1} with the prompt: {num_tokens_from_string(prompt)} tokens")
#         print("-------------------------------------------------------------------------------------------------------")
#     # after we have generated the summaries of each chunk of text, and appended them to the single summary string,
#     # we pass it into the final summary prompt
#     final_summary_prompt = f"""\n\nHuman: You will be given a set of summaries from a document. Create a cohesive 
#     summary from the provided individual summaries. The summary should very detailed and at least 10 pages. 
#     Summaries: {summary}
#             \n\nAssistant:"""
#     # print the total number of tokens being passed into the final summarization prompt
#     print(f"Number of tokens for this Chunk with the final prompt: {num_tokens_from_string(final_summary_prompt)}")
#     # generating the final summary of all the summaries we have previously generated.
#     return summarizer(final_summary_prompt)



def get_object_with_retry(bucket, key):
    max_retries=5
    retries = 0   
    backoff_base = 2
    max_backoff = 3  # Maximum backoff time in seconds
    s3 = boto3.client('s3')
    while retries < max_retries:
        try:
            response = s3.get_object(Bucket=bucket, Key=key)
            return response
        except ClientError as e:
            error_code = e.response['Error']['Code']
            if error_code == 'DecryptionFailureException':
                sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                print(f"Decryption failed, retrying in {sleep_time} seconds...")                
                time.sleep(sleep_time)               
                retries += 1
            elif e.response['Error']['Code'] == 'ModelStreamErrorException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
            else:
                raise e

    # If we reach this point, it means the maximum number of retries has been exceeded
    raise Exception(f"Failed to get object {key} from bucket {bucket} after {max_retries} retries.")


def save_chat_local(file_path, new_data, session_id):
    """Store long term chat history Local Disk"""   
    try:
        # Read the existing JSON data from the file
        with open(file_path, "r",encoding='utf-8') as file:
            existing_data = json.load(file)
        if session_id not in existing_data:
            existing_data[session_id]=[]
    except FileNotFoundError:
        # If the file doesn't exist, initialize an empty list
        existing_data = {session_id:[]}
    # Append the new data to the existing list
    from decimal import Decimal
    data = [{k: float(v) if isinstance(v, Decimal) else v for k, v in item.items()} for item in new_data]
    existing_data[session_id].extend(data)

    # Write the updated list back to the JSON file
    with open(file_path, "w",encoding="utf-8") as file:
        json.dump(existing_data, file)
        
def load_chat_local(file_path,session_id):
    """Load long term chat history from Local"""   
    try:
        # Read the existing JSON data from the file
        with open(file_path, "r",encoding='utf-8') as file:
            existing_data = json.load(file)
            if session_id in existing_data:
                existing_data=existing_data[session_id]
            else:
                existing_data=[]
    except FileNotFoundError:
        # If the file doesn't exist, initialize an empty list
        existing_data = []
    return existing_data
    
    
def process_files(files):

    result_string=""
    errors = []
    future_proxy_mapping = {} 
    futures = []

    with concurrent.futures.ProcessPoolExecutor() as executor:
        # Partial function to pass the handle_doc_upload_or_s3 function
        func = partial(handle_doc_upload_or_s3)   
        for file in files:
            future = executor.submit(func, file)
            future_proxy_mapping[future] = file
            futures.append(future)

        # Collect the results and handle exceptions
        for future in concurrent.futures.as_completed(futures):        
            file_url= future_proxy_mapping[future]
            try:
                result = future.result()               
                doc_name=os.path.basename(file_url)
                
                result_string+=f"<{doc_name}>\n{result}\n</{doc_name}>\n"
            except Exception as e:
                # Get the original function arguments from the Future object
                error = {'file': file_url, 'error': str(e)}
                errors.append(error)

    return errors, result_string

def handle_doc_upload_or_s3(file):
    """Handle various document format"""
    dir_name, ext = os.path.splitext(file)
    if  ext.lower() in [".pdf", ".png", ".jpg",".tif",".jpeg"]:   
        content=exract_pdf_text_aws(file)
    elif ".csv"  == ext.lower():
        content=parse_csv_from_s3(file)
    elif ext.lower() in [".xlsx", ".xls"]:
        content=table_parser_utills(file)   
    elif  ".json"==ext.lower():      
        obj=get_s3_obj_from_bucket_(file)
        content = json.loads(obj['Body'].read())  
    elif  ext.lower() in [".txt",".py"]:       
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
    elif ".docx" == ext.lower():       
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        docx_buffer = io.BytesIO(content)
        content = extract_text_and_tables(docx_buffer)
    elif ".pptx" == ext.lower():       
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        docx_buffer = io.BytesIO(content)        
        content = extract_text_from_pptx_s3(docx_buffer)
    else:            
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        doc_buffer = io.BytesIO(doc_content)
        content = textract.process(doc_buffer).decode()
    # Implement any other file extension logic 
    return content

class InvalidContentError(Exception):
    pass

def detect_encoding(s3_uri):
    """detect csv encoding"""
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", s3_uri)
    if match:
        bucket_name = match.group(1)
        key = match.group(2) 
    response = s3.get_object(Bucket=bucket_name, Key=key)
    content = response['Body'].read()
    result = chardet.detect(content)
    return result['encoding']

def parse_csv_from_s3(s3_uri):
    """read csv files"""
    try:
        # Detect the file encoding using chardet
        encoding = detect_encoding(s3_uri)        
        # Sniff the delimiter and read the CSV file
        df = pd.read_csv(s3_uri, delimiter=None, engine='python', encoding=encoding)
        return df.to_csv(index=False, sep=CSV_SEPERATOR)
    except Exception as e:
        raise InvalidContentError(f"Error: {e}")
    
def iter_block_items(parent):
    if isinstance(parent, Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("something's not right")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, parent)

def extract_text_and_tables(docx_path):
    """ Extract text from docx files"""
    document = DocxDocument(docx_path)
    content = ""
    current_section = ""
    section_type = None
    for block in iter_block_items(document):
        if isinstance(block, Paragraph):
            if block.text:
                if block.style.name == 'Heading 1':
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                        section_type = None  
                    section_type ="h1"
                    content += f"<{section_type}>{block.text}</{section_type}>\n"
                elif block.style.name== 'Heading 3':
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                    section_type = "h3"  
                    content += f"<{section_type}>{block.text}</{section_type}>\n"
                
                elif block.style.name == 'List Paragraph':
                    # Add to the current list section
                    if section_type != "list":
                        # Close the current section if it exists
                        if current_section:
                            content += f"{current_section}</{section_type}>\n"
                        section_type = "list"
                        current_section = "<list>"
                    current_section += f"{block.text}\n"
                elif block.style.name.startswith('toc'):
                    # Add to the current toc section
                    if section_type != "toc":
                        # Close the current section if it exists
                        if current_section:
                            content += f"{current_section}</{section_type}>\n"
                        section_type = "toc"
                        current_section = "<toc>"
                    current_section += f"{block.text}\n"
                else:
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                        section_type = None
                    
                    # Append the passage text without tagging
                    content += f"{block.text}\n"
        
        elif isinstance(block, DocxTable):
            # Add the current section before the table
            if current_section:
                content += f"{current_section}</{section_type}>\n"
                current_section = ""
                section_type = None

            content += "<table>\n"
            for row in block.rows:
                row_content = []
                for cell in row.cells:
                    cell_content = []
                    for nested_block in iter_block_items(cell):
                        if isinstance(nested_block, Paragraph):
                            cell_content.append(nested_block.text)
                        elif isinstance(nested_block, DocxTable):
                            nested_table_content = parse_nested_table(nested_block)
                            cell_content.append(nested_table_content)
                    row_content.append(CSV_SEPERATOR.join(cell_content))
                content += CSV_SEPERATOR.join(row_content) + "\n"
            content += "</table>\n"

    # Add the final section
    if current_section:
        content += f"{current_section}</{section_type}>\n"

    return content

def parse_nested_table(table):
    nested_table_content = "<table>\n"
    for row in table.rows:
        row_content = []
        for cell in row.cells:
            cell_content = []
            for nested_block in iter_block_items(cell):
                if isinstance(nested_block, Paragraph):
                    cell_content.append(nested_block.text)
                elif isinstance(nested_block, DocxTable):
                    nested_table_content += parse_nested_table(nested_block)
            row_content.append(CSV_SEPERATOR.join(cell_content))
        nested_table_content += CSV_SEPERATOR.join(row_content) + "\n"
    nested_table_content += "</table>"
    return nested_table_content



def extract_text_from_pptx_s3(pptx_buffer):
    """ Extract Text from pptx files"""
    presentation = Presentation(pptx_buffer)    
    text_content = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        text_content.append('\n'.join(slide_text))    
    return '\n\n'.join(text_content)
    
def exract_pdf_text_aws(file):    
    file_base_name=os.path.basename(file)
    dir_name, ext = os.path.splitext(file)
    # Checking if extracted doc content is in S3
    if USE_TEXTRACT:        
        if [x for x in get_s3_keys(f"{TEXTRACT_RESULT_CACHE_PATH}/") if file_base_name in x]:    
            response = get_object_with_retry(BUCKET, f"{TEXTRACT_RESULT_CACHE_PATH}/{file_base_name}.txt")
            text = response['Body'].read().decode()
            return text
        else:
            
            extractor = Textractor(region_name="us-east-1")
            # Asynchronous call, you will experience some wait time. Try caching results for better experience
            if "pdf" in ext:
                st.write("Asynchronous call, you may experience some wait time.")
                document = extractor.start_document_analysis(
                file_source=file,
                features=[TextractFeatures.LAYOUT,TextractFeatures.TABLES],       
                save_image=False,   
                s3_output_path=f"s3://{BUCKET}/textract_output/"
            )
            #Synchronous call
            else:
                document = extractor.analyze_document(
                file_source=file,
                features=[TextractFeatures.LAYOUT,TextractFeatures.TABLES],  
                save_image=False,
            )
            config = TextLinearizationConfig(
            hide_figure_layout=False,   
            hide_header_layout=False,    
            table_prefix="<table>",
            table_suffix="</table>",
            )
            # Upload extracted content to s3
            S3.put_object(Body=document.get_text(config=config), Bucket=BUCKET, Key=f"{TEXTRACT_RESULT_CACHE_PATH}/{file_base_name}.txt") 
            return document.get_text(config=config)
    else:
        s3=boto3.resource("s3")
        match = re.match("s3://(.+?)/(.+)", file)
        if match:
            bucket_name = match.group(1)
            key = match.group(2)
        if "pdf" in ext:            
            pdf_bytes = io.BytesIO()            
            s3.Bucket(bucket_name).download_fileobj(key, pdf_bytes)
            # Read the PDF from the BytesIO object
            pdf_bytes.seek(0)                      
            # Create a PDF reader object
            pdf_reader = PyPDF2.PdfReader(pdf_bytes)
            # Get the number of pages in the PDF
            num_pages = len(pdf_reader.pages)
            # Extract text from each page
            text = ''
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text += page.extract_text()
        else:
            img_bytes = io.BytesIO()
            s3.Bucket(bucket_name).download_fileobj(key, img_bytes)
            img_bytes.seek(0)
            image = Image.open(img_bytes)
            text = pytesseract.image_to_string(image)
        return text    

def strip_newline(cell):
    return str(cell).strip()

def table_parser_openpyxl(file):
    # Read from S3
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)    
        # Read Excel file from S3 into a buffer
        xlsx_buffer = io.BytesIO(obj['Body'].read())
        xlsx_buffer.seek(0)    
        # Load workbook
        wb = openpyxl.load_workbook(xlsx_buffer)    
        all_sheets_string=""
        # Iterate over each sheet in the workbook
        for sheet_name in wb.sheetnames:
            # all_sheets_name.append(sheet_name)
            worksheet = wb[sheet_name]

            all_merged_cell_ranges: list[CellRange] = list(
                worksheet.merged_cells.ranges
            )
            for merged_cell_range in all_merged_cell_ranges:
                merged_cell: Cell = merged_cell_range.start_cell
                worksheet.unmerge_cells(range_string=merged_cell_range.coord)
                for row_index, col_index in merged_cell_range.cells:
                    cell: Cell = worksheet.cell(row=row_index, column=col_index)
                    cell.value = merged_cell.value        
            # Convert sheet data to a DataFrame
            df = pd.DataFrame(worksheet.values)
            print(df.head(2))
            df = df.applymap(strip_newline)
            # Convert to string and tag by sheet name
            tabb=df.to_csv(sep=CSV_SEPERATOR, index=False, header=0)
            all_sheets_string+=f'<{sheet_name}>\n{tabb}\n</{sheet_name}>\n'
        return all_sheets_string
    else:
        raise Exception(f"{file} not formatted as an S3 path")

def calamaine_excel_engine(file):
    # # Read from S3
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)    
        # Read Excel file from S3 into a buffer
        xlsx_buffer = io.BytesIO(obj['Body'].read())
        xlsx_buffer.seek(0)    
        all_sheets_string=""
        # Load the Excel file
        workbook = CalamineWorkbook.from_filelike(xlsx_buffer)
        # Iterate over each sheet in the workbook
        for sheet_name in workbook.sheet_names:
            # Get the sheet by name
            sheet = workbook.get_sheet_by_name(sheet_name)
            df = pd.DataFrame(sheet.to_python(skip_empty_area=False))
            df = df.applymap(strip_newline)
            tabb=df.to_csv(sep=CSV_SEPERATOR, index=False, header=0)
            all_sheets_string+=f'<{sheet_name}>\n{tabb}\n</{sheet_name}>\n'
        return all_sheets_string
    else:
        raise Exception(f"{file} not formatted as an S3 path")

def table_parser_utills(file):
    try:
        response= table_parser_openpyxl(file)
        if response:
            return response
        else:
            return calamaine_excel_engine(file)        
    except Exception as e:
        try:
            return calamaine_excel_engine(file)
        except Exception as e:
            raise Exception(str(e))


def put_db(params,messages):
    """Store long term chat history in DynamoDB"""    
    chat_item = {
        "UserId": st.session_state['userid'], # user id
        "SessionId": params["session_id"], # User session id
        "messages": [messages],  # 'messages' is a list of dictionaries
        "time":messages['time']
    }

    existing_item = DYNAMODB.Table(DYNAMODB_TABLE).get_item(Key={"UserId": st.session_state['userid'], "SessionId":params["session_id"]})
    if "Item" in existing_item:
        existing_messages = existing_item["Item"]["messages"]
        chat_item["messages"] = existing_messages + [messages]
    response = DYNAMODB.Table(DYNAMODB_TABLE).put_item(
        Item=chat_item
    )
    
    
def get_chat_history_db(params,cutoff,claude3):
    current_chat, chat_hist=[],[]
    if params['chat_histories']: 
        chat_hist=params['chat_histories'][-cutoff:]              
        for d in chat_hist:
            if d['image'] and claude3 and LOAD_DOC_IN_ALL_CHAT_CONVO:
                content=[]
                for img in d['image']:
                    s3 = boto3.client('s3')
                    match = re.match("s3://(.+?)/(.+)", img)
                    image_name=os.path.basename(img)
                    _,ext=os.path.splitext(image_name)
                    if "jpg" in ext: ext=".jpeg"                        
                    if match:
                        bucket_name = match.group(1)
                        key = match.group(2)    
                        obj = s3.get_object(Bucket=bucket_name, Key=key)
                        base_64_encoded_data = base64.b64encode(obj['Body'].read())
                        base64_string = base_64_encoded_data.decode('utf-8')                        
                    content.extend([{"type":"text","text":image_name},{
                      "type": "image",
                      "source": {
                        "type": "base64",
                        "media_type": f"image/{ext.lower().replace('.','')}",
                        "data": base64_string
                      }
                    }])
                content.extend([{"type":"text","text":d['user']}])
                current_chat.append({'role': 'user', 'content': content})
            elif d['document'] and LOAD_DOC_IN_ALL_CHAT_CONVO:
                doc='Here are the documents:\n'
                for docs in d['document']:
                    uploads=handle_doc_upload_or_s3(docs)
                    doc_name=os.path.basename(docs)
                    doc+=f"<{doc_name}>\n{uploads}\n</{doc_name}>\n"
                if not claude3 and d["image"]:
                    for docs in d['image']:
                        uploads=handle_doc_upload_or_s3(docs)
                        doc_name=os.path.basename(docs)
                        doc+=f"<{doc_name}>\n{uploads}\n</{doc_name}>\n"
                current_chat.append({'role': 'user', 'content': [{"type":"text","text":doc+d['user']}]})
            else:
                current_chat.append({'role': 'user', 'content': [{"type":"text","text":d['user']}]})
            current_chat.append({'role': 'assistant', 'content': d['assistant']})  
    else:
        chat_hist=[]
    return current_chat, chat_hist

  
def get_s3_keys(prefix):
    """list all keys in an s3 path"""
    s3 = boto3.client('s3')
    keys = []
    next_token = None
    while True:
        if next_token:
            response = s3.list_objects_v2(Bucket=BUCKET, Prefix=prefix, ContinuationToken=next_token)
        else:
            response = s3.list_objects_v2(Bucket=BUCKET, Prefix=prefix)
        if "Contents" in response:
            for obj in response['Contents']:
                key = obj['Key']
                name = key[len(prefix):]
                keys.append(name)
        if "NextContinuationToken" in response:
            next_token = response["NextContinuationToken"]
        else:
            break
    return keys

def get_s3_obj_from_bucket_(file):
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)    
        obj = s3.get_object(Bucket=bucket_name, Key=key)  
    return obj

def put_obj_in_s3_bucket_(docs):
    file_name=os.path.basename(docs.name)
    file_path=f"{S3_DOC_CACHE_PATH}/{file_name}"
    S3.put_object(Body=docs.read(),Bucket= BUCKET, Key=file_path)
    return f"s3://{BUCKET}/{file_path}"


def bedrock_streemer(params,response, handler):
    stream = response.get('body')
    answer = ""    
    if stream:
        for event in stream:
            chunk = event.get('chunk')
            if  chunk:
                chunk_obj = json.loads(chunk.get('bytes').decode())
                if "delta" in chunk_obj:                    
                    delta = chunk_obj['delta']
                    if "text" in delta:
                        text=delta['text'] 
                        # st.write(text, end="")                        
                        answer+=str(text)       
                        handler.markdown(answer.replace("$","USD ").replace("%", " percent"))
                        
                if "amazon-bedrock-invocationMetrics" in chunk_obj:
                    st.session_state['input_token'] = chunk_obj['amazon-bedrock-invocationMetrics']['inputTokenCount']
                    st.session_state['output_token'] =chunk_obj['amazon-bedrock-invocationMetrics']['outputTokenCount']
                    pricing=st.session_state['input_token']*pricing_file[f"anthropic.{params['model']}"]["input"]+st.session_state['output_token'] *pricing_file[f"anthropic.{params['model']}"]["output"]
                    st.session_state['cost']+=pricing             
    return answer

def bedrock_claude_(params,chat_history,system_message, prompt,model_id,image_path=None, handler=None):
    content=[]
    if image_path:       
        if not isinstance(image_path, list):
            image_path=[image_path]      
        for img in image_path:
            s3 = boto3.client('s3')
            match = re.match("s3://(.+?)/(.+)", img)
            image_name=os.path.basename(img)
            _,ext=os.path.splitext(image_name)
            if "jpg" in ext.lower(): ext=".jpeg"                        
            if match:
                bucket_name = match.group(1)
                key = match.group(2)    
                obj = s3.get_object(Bucket=bucket_name, Key=key)
                base_64_encoded_data = base64.b64encode(obj['Body'].read())
                base64_string = base_64_encoded_data.decode('utf-8')
            content.extend([{"type":"text","text":image_name},{
              "type": "image",
              "source": {
                "type": "base64",
                "media_type": f"image/{ext.lower().replace('.','')}",
                "data": base64_string
              }
            }])
    content.append({
        "type": "text",
        "text": prompt
            })
    chat_history.append({"role": "user",
            "content": content})
    # print(system_message)
    prompt = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 1500,
        "temperature": 0,
        "system":system_message,
        "messages": chat_history
    }
    
    prompt = json.dumps(prompt)
    response = bedrock_runtime.invoke_model_with_response_stream(body=prompt, modelId=model_id, accept="application/json", contentType="application/json")
    answer=bedrock_streemer(params,response, handler) 
    return answer

def _invoke_bedrock_with_retries(params,current_chat, chat_template, question, model_id, image_path, handler):
    max_retries = 10
    backoff_base = 2
    max_backoff = 3  # Maximum backoff time in seconds
    retries = 0

    while True:
        try:
            response = bedrock_claude_(params,current_chat, chat_template, question, model_id, image_path, handler)
            return response
        except ClientError as e:
            if e.response['Error']['Code'] == 'ThrottlingException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            elif e.response['Error']['Code'] == 'ModelStreamErrorException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            elif e.response['Error']['Code'] == 'EventStreamError':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            else:
                # Some other API error, rethrow
                raise
                

def get_session_ids_by_user(table_name, user_id):
    """
    Get Session Ids and corresponding top message for a user to populate the chat history drop down on the front end
    """
    if DYNAMODB_TABLE:
        table = DYNAMODB.Table(table_name)
        message_list={}
        session_ids = []
        args = {
            'KeyConditionExpression': Key('UserId').eq(user_id)
        }
        while True:
            response = table.query(**args)
            session_ids.extend([item['SessionId'] for item in response['Items']])
            if 'LastEvaluatedKey' not in response:
                break
            args['ExclusiveStartKey'] = response['LastEvaluatedKey']

        for session_id in session_ids:
            try:
                message_list[session_id]=DYNAMODB.Table(table_name).get_item(Key={"UserId": user_id, "SessionId":session_id})['Item']['messages'][0]['user']
            except Exception as e:
                print(e)
                pass
    else:
        try:
            message_list={}
            # Read the existing JSON data from the file
            with open(LOCAL_CHAT_FILE_NAME, "r", encoding='utf-8') as file:
                existing_data = json.load(file)
            for session_id in existing_data:
                message_list[session_id]=existing_data[session_id][0]['user']
            
        except FileNotFoundError:
            # If the file doesn't exist, initialize an empty list
            message_list = {}
    return message_list


# def query_llm(params, handler):
#     """
#     Function takes a user query and a uploaded document. Caches documents in S3
#     passing a document is optional
#     """  

#     if not isinstance(params['upload_doc'], list):
#         raise TypeError("documents must be in a list format")        
#     # Check if Claude3 model is used and handle images with the CLAUDE3 Model
#     claude3=False
#     model='anthropic.'+params['model']
#     if "sonnet" in model or "haiku" in model:
#         model+="-20240229-v1:0" if "sonnet" in model else "-20240307-v1:0"
#         claude3=True
#     # Retrieve past chat history   
#     current_chat,chat_hist=get_chat_history_db(params, CHAT_HISTORY_LENGTH,claude3)

#     ## prompt template for when a user uploads a doc
#     doc_path=[]
#     image_path=[]
#     full_doc_path=[]
#     doc=""
#     if params['upload_doc']:  
#         doc='I have provided documents and/or images.\n'
#         for ids,docs in enumerate(params['upload_doc']):
#             file_name=docs.name
#             _,extensions=os.path.splitext(file_name)
#             docs=put_obj_in_s3_bucket_(docs)
#             full_doc_path.append(docs)
#             if extensions.lower() in [".jpg",".jpeg",".png",".gif",".webp"] and claude3:
#                 image_path.append(docs)
#                 continue
        
#         doc_path = [item for item in full_doc_path if item not in image_path]
#         errors, result_string=process_files(doc_path)    
#         if errors:
#             st.error(errors)
#         print(len(result_string))

#         if len(result_string) > 175000:
#             st.error(f"Content exceeds the limit of 200k tokens")

#         doc+= result_string
#         with open("prompt/doc_chat.txt","r", encoding="utf-8") as f:
#             chat_template=f.read()  
#     else:        
#         # Chat template for open ended query
#         with open("prompt/chat.txt","r",encoding="utf-8") as f:
#             chat_template=f.read()

#     response=_invoke_bedrock_with_retries(params,current_chat, chat_template, doc+params['question'], model, image_path, handler)
#     # log the following items to dynamodb
#     chat_history={"user":params['question'],
#     "assistant":response,
#     "image":image_path,
#     "document":doc_path,
#     "modelID":model,
#     "time":str(time.time()),
#     "input_token":round(st.session_state['input_token']) ,
#     "output_token":round(st.session_state['output_token'])} 
#     #store convsation memory and user other items in DynamoDB table
#     if DYNAMODB_TABLE:
#         put_db(params,chat_history)
#     # use local memory for storage
#     else:
#         save_chat_local(LOCAL_CHAT_FILE_NAME,[chat_history], params["session_id"])  
#     return response


def query_llm(params, handler):
    """
    Function takes a user query and a uploaded document. Caches documents in S3
    passing a document is optional
    """  

    if not isinstance(params['upload_doc'], list):
        raise TypeError("documents must be in a list format")        
    # Check if Claude3 model is used and handle images with the CLAUDE3 Model
    claude3=False
    model='anthropic.'+params['model']
    if "sonnet" in model or "haiku" in model:
        model+="-20240229-v1:0" if "sonnet" in model else "-20240307-v1:0"
        claude3=True
    # Retrieve past chat history   
    current_chat,chat_hist=get_chat_history_db(params, CHAT_HISTORY_LENGTH,claude3)

    ## prompt template for when a user uploads a doc
    doc_path=[]
    image_path=[]
    full_doc_path=[]
    doc=""
    response = ''
    if params['upload_doc']:  
        doc='I have provided documents and/or images.\n'
        for ids,docs in enumerate(params['upload_doc']):
            file_name=docs.name
            _,extensions=os.path.splitext(file_name)
            docs=put_obj_in_s3_bucket_(docs)
            full_doc_path.append(docs)
            if extensions.lower() in [".jpg",".jpeg",".png",".gif",".webp"] and claude3:
                image_path.append(docs)
                continue
        
        doc_path = [item for item in full_doc_path if item not in image_path]
        errors, result_string=process_files(doc_path)    
        if errors:
            st.error(errors)
        print(len(result_string))
        
        with open("prompt/doc_chat.txt","r", encoding="utf-8") as f:
                chat_template=f.read() 

        if len(result_string) > 175000:
            st.error(f"Content exceeds the limit of 200k tokens. Working on it...")
            if DOCUMENT_SUMMARIZATION:
                doc+= Chunk_and_Summarize(result_string)
                response = doc
            elif QNA:
                doc+= InMemoryFAISS(result_string, chat_hist, params['question'])
                print(f"QNA: {doc}" )
                #send to LLM

        # if len(result_string):
        #     # st.error(f"Content exceeds the limit of 200k tokens. Working on it...")
        #     if DOCUMENT_SUMMARIZATION:
        #         doc+= Chunk_and_Summarize(result_string)
        #         response = doc
        #     elif QNA:
        #         doc+= InMemoryFAISS(result_string, chat_hist, params['question'])
        #         print(f"QNA: {doc}" )
        #         response = doc
        #         #send to LLM 
 
        else:
            doc+= result_string            
            response=_invoke_bedrock_with_retries(params,current_chat, chat_template, doc+params['question'], model, image_path, handler)


        

    else:        
        # Chat template for open ended query
        with open("prompt/chat.txt","r",encoding="utf-8") as f:
            chat_template=f.read()
        response=_invoke_bedrock_with_retries(params,current_chat, chat_template, doc+params['question'], model, image_path, handler)

    
    # log the following items to dynamodb
    chat_history={"user":params['question'],
    "assistant":response,
    "image":image_path,
    "document":doc_path,
    "modelID":model,
    "time":str(time.time()),
    "input_token":round(st.session_state['input_token']) ,
    "output_token":round(st.session_state['output_token'])} 
    #store convsation memory and user other items in DynamoDB table
    if DYNAMODB_TABLE:
        put_db(params,chat_history)
    # use local memory for storage
    else:
        save_chat_local(LOCAL_CHAT_FILE_NAME,[chat_history], params["session_id"])  
    return response


def get_chat_historie_for_streamlit(params):
    """
    This function retrieves chat history stored in a dynamoDB table partitioned by a userID and sorted by a SessionID
    """
    if DYNAMODB_TABLE:
        chat_histories = DYNAMODB.Table(DYNAMODB_TABLE).get_item(Key={"UserId": st.session_state['userid'], "SessionId":params["session_id"]})
        if "Item" in chat_histories:
            chat_histories=chat_histories['Item']['messages'] 
        else:
            chat_histories=[]
    else:
        chat_histories=load_chat_local(LOCAL_CHAT_FILE_NAME,params["session_id"])         

# Constructing the desired list of dictionaries
    formatted_data = []   
    if chat_histories:
        for entry in chat_histories:           
            image_files=[os.path.basename(x) for x in entry.get('image', [])]
            doc_files=[os.path.basename(x) for x in entry.get('document', [])]
            assistant_attachment = '\n\n'.join(image_files+doc_files)
            
            formatted_data.append({
                "role": "user",
                "content": entry["user"],
            })
            formatted_data.append({
                "role": "assistant",
                "content": entry["assistant"],
                "attachment": assistant_attachment
            })
    else:
        chat_histories=[]            
    return formatted_data,chat_histories



def get_key_from_value(dictionary, value):
    return next((key for key, val in dictionary.items() if val == value), None)
    
def chat_bedrock_(params):
    # st.title('ESG Data Assistant 2')
    params['chat_histories']=[]
    if params["session_id"].strip():
        st.session_state.messages, params['chat_histories']=get_chat_historie_for_streamlit(params)
    for message in st.session_state.messages:

        with st.chat_message(message["role"]):   
            if "```" in message["content"]:
                st.markdown(message["content"],unsafe_allow_html=True )
            else:
                st.markdown(message["content"].replace("$", "\$"),unsafe_allow_html=True )
            if message["role"]=="assistant":
                if message["attachment"]:
                    with st.expander(label="**attachments**"):
                        st.markdown( message["attachment"])
    if prompt := st.chat_input("Whats up?"):        
        st.session_state.messages.append({"role": "user", "content": prompt})        
        with st.chat_message("user"):             
            st.markdown(prompt.replace("$", "\$"),unsafe_allow_html=True )
        with st.chat_message("assistant"): 
            message_placeholder = st.empty()
            time_now=time.time()            
            params["question"]=prompt
            answer=query_llm(params, message_placeholder)
            print(f"from llm: {answer}")
            message_placeholder.markdown(answer.replace("$", "\$"),unsafe_allow_html=True )
            st.session_state.messages.append({"role": "assistant", "content": answer}) 
        st.rerun()
        
def app_sidebar():
    with st.sidebar:   
        st.metric(label="Session Cost", value=f"${round(st.session_state['cost'],3)}") 
        st.markdown(f"**{st.session_state['firstname']} {st.session_state['lastname']}**")
        st.write("-----")
        button=st.button("New Chat", type ="primary")
        models=[ 'claude-3-sonnet','claude-3-haiku','claude-instant-v1','claude-v2:1', 'claude-v2']
        model=st.selectbox('**Model**', models)
        params={"model":model} 
        user_sess_id=get_session_ids_by_user(DYNAMODB_TABLE, st.session_state['userid'])
        float_keys = {float(key): value for key, value in user_sess_id.items()}
        sorted_messages = sorted(float_keys.items(), reverse=True)      
        sorted_messages.insert(0, (float(st.session_state['user_sess']),"New Chat"))        
        if button:
            st.session_state['user_sess'] = str(time.time())
            sorted_messages.insert(0, (float(st.session_state['user_sess']),"New Chat"))      
        st.session_state['chat_session_list'] = dict(sorted_messages)
        chat_items=st.selectbox("**Chat Sessions**",st.session_state['chat_session_list'].values(),key="chat_sessions")
        session_id=get_key_from_value(st.session_state['chat_session_list'], chat_items)   
        file = st.file_uploader('Upload a document', accept_multiple_files=True, help="pdf,csv,txt,png,jpg,xlsx,json,py doc format supported") 
        if file and LOAD_DOC_IN_ALL_CHAT_CONVO:
            st.warning('You have set **load-doc-in-chat-history** to true. For better performance, remove upload files (by clicking **X**) **AFTER** first query **RESPONSE** on uploaded files. See the README for more info', icon="⚠️")
        params={"model":model, "session_id":str(session_id), "chat_item":chat_items, "upload_doc":file }    
        st.session_state['count']=1
        return params

# Get the Flask app's authentication endpoint URL from the environment variable
AUTH_ENDPOINT = os.environ.get("REACT_APP_USERS_SERVICE_URL", "") + "/auth/status"
print(f"AUTH_ENDPOINT: {AUTH_ENDPOINT}")

def get_token_from_local_storage(k):

    v = st_javascript(
        f"localStorage.getItem('{k}');"
    )
    token = v
    print(f"Token: {token}")
    return token

def is_authenticated():
    """
    Check if the user is authenticated by sending a request to the authentication endpoint.
    Returns True if the user is authenticated, False otherwise.
    """
    token = get_token_from_local_storage("authToken")
    # s = requests.get("/users/ping")
    # ping_url = os.environ.get("REACT_APP_USERS_SERVICE_URL", "") + "/users/ping"
    # ping_url = f"http://users:5000/users/ping"
    # s = requests.get(ping_url)
    # print(f"Ping: {s.json()}")

    try:        
        headers = {"Authorization": f"Bearer {token}"}
        # print(f"headers: {headers}")
        # print(f"AUTH_ENDPOINT: {AUTH_ENDPOINT}")
        #response = requests.get(AUTH_ENDPOINT, headers=headers)
        response = requests.get(f"http://users:5000/auth/status", headers=headers)
        # response = requests.get(f"http://localhost/auth/status", headers=headers)
        # time.sleep(3)
        print(f"Response: {response.json()}")
        response.raise_for_status()
        data = response.json()

        if data["status"] == "success":
            if "firstname" not in st.session_state:
                st.session_state['firstname']= data["data"]["firstname"]
            if "lastname" not in st.session_state:
                st.session_state['lastname']= data["data"]["lastname"]
            config_file["UserId"] = data["data"]["email"]
            # if "userid" not in st.session_state:
            st.session_state['userid']= config_file["UserId"]
            return True
    except requests.exceptions.RequestException:
        pass
    return False


def main():
    
    st.markdown("**ESD Data Assistant 2**")
    st.header("Document Analysis")
    # st.markdown(r"$\textsf{\small Enter text here}$")
    
    description = """
    The app is designed to leverage the power of Generative AI (GenAI) to assist users in analyzing and gaining insights from ESG (Environmental, Social, and Governance) sustainability and climate change-related documents. Users can upload various types of documents, such as reports, research papers, policy documents, and more, and then perform a range of tasks using GenAI.
    """

    st.markdown(description)

    features = [
        "**Document Summarization:** Generate concise summaries of lengthy documents, highlighting the key points and main takeaways related to ESG sustainability and climate change.",
        "**Topic Extraction:** Identify and extract relevant topics, themes, and key concepts from the uploaded documents, enabling users to quickly grasp the main focus areas.",
        "**Sentiment Analysis:** Analyze the sentiment expressed in the documents regarding specific ESG initiatives, climate change policies, or corporate sustainability efforts.",
        "**Question Answering:** Ask natural language questions related to the content of the documents, and receive precise answers generated by the GenAI model.",
        "**Data Extraction:** Extract specific data points, statistics, and numerical values from the documents, such as emission levels, carbon footprint metrics, or ESG performance indicators.",
        "**Trend Analysis:** Identify emerging trends, patterns, and potential future developments related to ESG sustainability and climate change based on the document content.",
        "**Comparative Analysis:** Compare multiple documents side-by-side, highlighting similarities, differences, and contrasting viewpoints on specific ESG or climate change topics.",
        "**Report Generation:** Generate comprehensive reports summarizing the key findings, insights, and recommendations based on the document analysis, tailored for different stakeholders (e.g., policymakers, corporate leaders, investors).",
        "**Language Translation:** Translate documents or specific sections into multiple languages, enabling broader accessibility and understanding of ESG and climate change information.",
        "**Customizable Prompts:** Allow users to provide custom prompts or instructions to the GenAI model, enabling more specific and targeted analysis based on their unique requirements."
    ]

    cols = st.columns(2)
    for i, feature in enumerate(features):
        cols[i % 2].markdown(f"- {feature}")

    st.markdown("""
    The app aims to streamline the analysis of ESG sustainability and climate change-related documents, empowering users with the ability to quickly extract valuable insights, identify trends, and make data-driven decisions. By leveraging the power of GenAI, users can save time and effort while gaining a deeper understanding of these critical topics.
    """)


    if is_authenticated():
        params=app_sidebar()
        chat_bedrock_(params)
    else:
        login_url = urljoin(os.environ.get("REACT_APP_USERS_SERVICE_URL", ""), "/login")
        print(login_url)
        st.markdown(f"You are not authenticated. Please [log in]({login_url}) to access the app.")

 
    
if __name__ == '__main__':
    main() 